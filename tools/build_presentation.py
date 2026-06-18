from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Mentoria_Compass_Mogger_Presentation.pptx"
LINKS = ROOT / "Mentoria_Compass_Mogger_Links.txt"

PRODUCT_URL = "https://mentoria-compass.vercel.app"
GITHUB_URL = "https://github.com/askarshabdanov1-collab/Mentoria"


slides = [
    {
        "label": "Проект",
        "title": "Mentoria Compass",
        "subtitle": "Full-stack EdTech-платформа для возможностей, асинхронного обучения, AI-roadmap и анализа CV",
        "bullets": [
            "Команда: Mogger",
            "Капитан: Шабданов Аскар / @LawDirr",
            f"Готовый продукт: {PRODUCT_URL}",
            f"GitHub-репозиторий: {GITHUB_URL}",
            "Демо-аккаунты: ученик amina@mentoria.demo / demo123, админ admin@mentoria.demo / admin123",
        ],
    },
    {
        "label": "Проблема",
        "title": "Зачем Mentoria нужна платформа",
        "subtitle": "Telegram и живые занятия перестают масштабироваться, когда организация растёт",
        "bullets": [
            "Ученики пропускают живые уроки из-за школы, экзаменов, часовых поясов, интернета и личного расписания.",
            "Возможности разбросаны по сайтам, Telegram-каналам, чатам и документам.",
            "Ученикам сложно понять, какие конкурсы, стипендии и программы подходят их классу, интересам и целям.",
            "Администраторам Mentoria нужно добавлять курсы и возможности без пересборки сайта и ручных постов.",
            "Партнёрам, школам и спонсорам важно видеть профессиональный масштабируемый продукт, а не только чат.",
        ],
    },
    {
        "label": "Аудитория",
        "title": "Для кого создан продукт",
        "subtitle": "Ученики 8-11 классов, которые хотят развиваться за пределами школы",
        "bullets": [
            "Ученики из Казахстана и других стран, которые ищут академические и внеклассные возможности.",
            "Ученики, готовящиеся к поступлению, IELTS/SAT, стипендиям и международным программам.",
            "Ученики, которым интересны STEM, бизнес, финансы, social impact, программирование, наука и английский.",
            "Администраторы и менторы Mentoria, которым нужно управлять контентом и прогрессом учеников в одном месте.",
        ],
    },
    {
        "label": "Решение",
        "title": "Что делает Mentoria Compass",
        "subtitle": "Один продукт, который соединяет поиск возможностей, обучение, прогресс и заявки",
        "bullets": [
            "Регистрация и личный кабинет ученика с сохранёнными возможностями, прогрессом и дедлайнами.",
            "Каталог возможностей с поиском, фильтрами и match score на основе профиля.",
            "Асинхронные курсы с уроками, видео-плейсхолдерами, мини-заданиями и отслеживанием прогресса.",
            "AI Analysis: readiness score, риски, сильные стороны и недельный план действий.",
            "Админ-панель для публикации возможностей и курсов через backend API.",
        ],
    },
    {
        "label": "Сценарий",
        "title": "Основной пользовательский путь",
        "subtitle": "Сценарий удобно показать в 4-минутном демо",
        "bullets": [
            "Ученик регистрируется или входит и заполняет профиль: класс, страна, школа, английский, интересы и цели.",
            "Compass dashboard рекомендует возможности и курсы с match score.",
            "Ученик сохраняет возможность, видит дедлайн и начинает подходящий курс.",
            "Ученик завершает урок, а личный кабинет автоматически обновляет прогресс.",
            "Ученик открывает AI Analysis и получает следующие шаги для заявок и обучения.",
            "Админ входит и добавляет новую возможность, которая сразу появляется в каталоге.",
        ],
    },
    {
        "label": "AI",
        "title": "AI-roadmap и рекомендации",
        "subtitle": "Практичный анализ, который работает сразу без внешних API-ключей",
        "bullets": [
            "Readiness score считается по полноте профиля, сохранённым возможностям, прогрессу курсов и срочности дедлайнов.",
            "Система находит риски: нет сохранённых возможностей, низкий прогресс, близкие дедлайны, слабый фокус на портфолио.",
            "Ученик получает конкретные следующие действия, а не общие советы.",
            "Недельный план превращает рекомендации в задачи: выбрать возможность, пройти урок, написать ответ, спросить ментора, отправить заявку.",
            "Это повышает удержание: ученик всегда понимает, что делать дальше.",
        ],
    },
    {
        "label": "CV Review",
        "title": "CV Review: найти дыры и улучшить",
        "subtitle": "Ученик может загрузить или вставить CV и получить фидбек под заявки",
        "bullets": [
            "Платформа проверяет, есть ли контакты, образование, проекты, навыки, достижения и ссылки.",
            "Система находит слабые места: нет цифр, пассивные bullets, нет portfolio/GitHub-ссылок, неясные достижения.",
            "Ученик получает CV readiness score и приоритетные дыры, которые нужно закрыть первыми.",
            "Инструмент предлагает конкретные улучшения: сильный summary, измеримый impact, проектный раздел и переписанные bullets.",
            "Анализы CV сохраняются в аккаунте ученика, чтобы можно было вернуться к прогрессу позже.",
        ],
    },
    {
        "label": "Технологии",
        "title": "Техническая архитектура",
        "subtitle": "Full-stack приложение, задеплоенное в production",
        "bullets": [
            "Frontend: React + Vite dashboard с адаптивной навигацией и продуктовым UI.",
            "Backend: Node API с auth, профилем, прогрессом, сохранёнными элементами, админ-контентом и CV analysis routes.",
            "Аутентификация: регистрация, вход, выход и bearer-token sessions.",
            "Deployment: production на Vercel, подключённый к GitHub-репозиторию.",
            "Текущее хранение: JSON seed/serverless storage для хакатона; следующий шаг для production — Supabase/Neon Postgres.",
        ],
    },
    {
        "label": "Impact",
        "title": "Влияние для Mentoria",
        "subtitle": "Как продукт помогает Mentoria масштабироваться за пределы Telegram",
        "bullets": [
            "Масштабирует асинхронное обучение: ученики могут учиться даже без live-занятий.",
            "Централизует возможности: курсы, дедлайны, рекомендации и сохранённые элементы находятся в одном месте.",
            "Повышает удержание: progress bars, дедлайны, AI next steps и CV feedback удерживают вовлечённость.",
            "Усиливает профессиональный имидж: Mentoria может показать школам, спонсорам и партнёрам реальную платформу.",
            "Снижает нагрузку на админов: возможности и курсы добавляются через панель, без пересборки и ручных постов.",
        ],
    },
    {
        "label": "Roadmap",
        "title": "Что дальше",
        "subtitle": "Путь от хакатонного MVP к реальному продукту Mentoria",
        "bullets": [
            "Подключить постоянную production-базу данных: Supabase или Neon Postgres.",
            "Добавить Telegram/email-напоминания по дедлайнам сохранённых возможностей.",
            "Добавить портал менторов для загрузки уроков, проверки заданий и фидбека.",
            "Добавить мультиязычный интерфейс: русский, английский и казахский.",
            "Добавить сертификаты, leaderboard и расширенный roadmap для 8, 9, 10 и 11 классов.",
            "Добавить реальную AI-интеграцию для глубокого анализа CV и эссе при наличии API-ключей.",
        ],
    },
]


def add_textbox(slide, x, y, w, h, text, size=24, color=RGBColor(246, 248, 239), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(2.35), Inches(11.45), Inches(4.62))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(7)
        p.font.size = Pt(17 if len(bullets) > 5 else 18)
        p.font.color.rgb = RGBColor(230, 234, 218)


def add_accent(slide, prs, index, label):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(10, 12, 11)

    top = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(223, 255, 85)
    top.line.fill.background()

    badge = slide.shapes.add_shape(1, Inches(10.35), Inches(0.42), Inches(2.1), Inches(0.48))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(90, 215, 255)
    badge.line.fill.background()
    tf = badge.text_frame
    tf.text = f"{index:02d} / {label}"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(13)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(10, 12, 11)

    footer = slide.shapes.add_textbox(Inches(0.75), Inches(7.05), Inches(11.7), Inches(0.25))
    footer.text_frame.text = f"{PRODUCT_URL}  |  {GITHUB_URL}"
    footer.text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)
    footer.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(115, 132, 114)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, item in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        add_accent(slide, prs, idx, item["label"])
        add_textbox(slide, 0.72, 0.58, 9.5, 0.44, "MENTORIA COMPASS / КОМАНДА MOGGER", 13, RGBColor(223, 255, 85), True)
        add_textbox(slide, 0.7, 0.98, 11.15, 0.82, item["title"], 38, RGBColor(255, 254, 240), True)
        add_textbox(slide, 0.73, 1.78, 11.2, 0.45, item["subtitle"], 17, RGBColor(170, 184, 169), False)
        add_bullets(slide, item["bullets"])

    prs.save(OUT)
    LINKS.write_text(
        "Mentoria Compass / Команда Mogger\n"
        "Капитан: Шабданов Аскар (@LawDirr)\n\n"
        f"Готовый продукт:\n{PRODUCT_URL}\n\n"
        f"GitHub-репозиторий:\n{GITHUB_URL}\n\n"
        "Демо-аккаунты:\n"
        "Ученик: amina@mentoria.demo / demo123\n"
        "Админ: admin@mentoria.demo / admin123\n",
        encoding="utf-8",
    )
    print(OUT)
    print(LINKS)


if __name__ == "__main__":
    main()
